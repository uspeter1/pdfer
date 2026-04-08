import io
import shutil
import subprocess
import tempfile
from pathlib import Path


def _libreoffice_convert(src: Path, out: Path, fmt: str) -> bool:
    """Convert via LibreOffice headless. Returns True on success."""
    lo = shutil.which('soffice') or shutil.which('libreoffice')
    if not lo:
        return False

    ext = fmt.lower()  # 'docx' or 'pptx'

    # PDF import requires explicit infilter so LO knows to treat the file as PDF
    if ext == 'docx':
        infilter = 'writer_pdf_import'
    else:
        infilter = 'impress_pdf_import'

    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        try:
            result = subprocess.run(
                [
                    lo, '--headless',
                    f'--infilter={infilter}',
                    '--convert-to', ext,
                    '--outdir', str(td_path),
                    str(src),
                ],
                capture_output=True,
                timeout=120,
            )
        except subprocess.TimeoutExpired:
            return False

        # LibreOffice writes {stem}.{ext} into outdir
        candidate = td_path / f'{src.stem}.{ext}'
        if result.returncode == 0 and candidate.exists() and candidate.stat().st_size > 0:
            shutil.move(str(candidate), str(out))
            return True

    return False


def _to_docx_fallback(src: Path, out: Path):
    """pdf2docx fallback (slow but layout-aware)."""
    try:
        from pdf2docx import Converter
    except ImportError:
        raise RuntimeError('pdf2docx is not installed. Run: pip install pdf2docx')
    cv = Converter(str(src))
    cv.convert(str(out), start=0, end=None)
    cv.close()


def _to_pptx_fallback(src: Path, out: Path):
    """python-pptx fallback: render each page as an image slide."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        raise RuntimeError('python-pptx is not installed. Run: pip install python-pptx')
    import fitz

    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(6858000)

    with fitz.open(str(src)) as doc:
        for page in doc:
            pw, ph = page.rect.width, page.rect.height
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes('png')

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            if pw > 0 and ph > 0:
                prs.slide_height = Emu(int(prs.slide_width * ph / pw))

            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left=0, top=0,
                width=prs.slide_width, height=prs.slide_height,
            )

    prs.save(str(out))


def run(input_files: list, params: dict, work_dir: Path) -> list:
    """Convert PDF to DOCX or PPTX."""
    fmt = str(params.get('format', 'DOCX')).upper()
    ext = fmt.lower()

    outputs = []
    for f in input_files:
        if f.suffix.lower() != '.pdf':
            outputs.append(f)
            continue

        out = work_dir / f'{f.stem}.{ext}'

        # Primary: LibreOffice (fast, good quality)
        if _libreoffice_convert(f, out, ext):
            outputs.append(out)
            continue

        # Fallback
        if fmt == 'DOCX':
            _to_docx_fallback(f, out)
        else:
            _to_pptx_fallback(f, out)

        outputs.append(out)

    return outputs
